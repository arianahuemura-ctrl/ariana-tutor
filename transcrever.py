import yt_dlp
import fitz
import requests
import os
import subprocess
import asyncio
import edge_tts
import time
from datetime import datetime
from pptx import Presentation
from groq import Groq
from ddgs import DDGS
from config import TELEGRAM_TOKEN, GROQ_API_KEY
from pydub import AudioSegment

TELEGRAM_URL = f"https://api.telegram.org/bot{TELEGRAM_TOKEN}"
CHAT_ID = "6968289835"
groq_client = Groq(api_key=GROQ_API_KEY)

COOKIES_FILE = "/home/ubuntu/ariana-tutor/cookies.txt"

QUERIES_FONTES = [
    "{tema} tutorial site:developer.mozilla.org OR site:web.dev OR site:devdocs.io",
    "{tema} UX design site:nngroup.com OR site:smashingmagazine.com OR site:lawsofux.com",
    "{tema} definition site:dictionary.cambridge.org OR site:merriam-webster.com",
    "{tema} best practices site:wikipedia.org OR site:roadmap.sh OR site:refactoring.guru",
]

def baixar_audio_youtube(url, tentativas=3):
    print(f"Baixando audio do YouTube: {url}")
    for tentativa in range(1, tentativas + 1):
        try:
            resultado = subprocess.run([
                'yt-dlp',
                '--js-runtimes', 'node',
                '--format', 'bestaudio/best',
                '--extract-audio',
                '--audio-format', 'mp3',
                '--output', '/tmp/aula.%(ext)s',
                '--quiet',
                '--socket-timeout', '60',
                '--retries', '5',
                '--cookies', COOKIES_FILE,
                url
            ], capture_output=True, text=True)

            if os.path.exists('/tmp/aula.mp3'):
                return '/tmp/aula.mp3'
            for f in os.listdir('/tmp'):
                if f.startswith('aula') and f.endswith('.mp3'):
                    return f'/tmp/{f}'
            if resultado.returncode != 0:
                raise Exception(resultado.stderr)
        except Exception as e:
            print(f"Tentativa {tentativa}/{tentativas} falhou: {e}")
            if tentativa < tentativas:
                time.sleep(10 * tentativa)
    print(f"AVISO: Pulando vídeo {url}")
    return None

def particionar_audio(caminho_audio, minutos=10):
    print(f"Particionando audio em blocos de {minutos} minutos...")
    audio = AudioSegment.from_mp3(caminho_audio)
    bloco_ms = minutos * 60 * 1000
    partes = []
    for i, inicio in enumerate(range(0, len(audio), bloco_ms)):
        parte = audio[inicio:inicio + bloco_ms]
        caminho = f"/tmp/parte_{i}.mp3"
        parte.export(caminho, format="mp3")
        partes.append(caminho)
    print(f"Audio dividido em {len(partes)} partes.")
    return partes

def transcrever_audio(caminho_audio):
    print("Transcrevendo audio com Whisper via Groq API...")
    partes = particionar_audio(caminho_audio)
    transcricao_completa = ""
    for i, parte in enumerate(partes):
        print(f"Transcrevendo parte {i+1}/{len(partes)}...")
        try:
            with open(parte, "rb") as f:
                resultado = groq_client.audio.transcriptions.create(
                    model="whisper-large-v3",
                    file=f,
                    language="pt"
                )
            transcricao_completa += f"\n[Parte {i+1}]\n{resultado.text}\n"
        except Exception as e:
            print(f"Erro na parte {i+1}: {e}")
            transcricao_completa += f"\n[Parte {i+1}]\n[Erro na transcrição desta parte]\n"
        finally:
            if os.path.exists(parte):
                os.remove(parte)
    return transcricao_completa

def extrair_texto_pdf(caminho_pdf):
    print("Lendo PDF...")
    doc = fitz.open(caminho_pdf)
    texto = ""
    for pagina in doc:
        texto += pagina.get_text()
    return texto

def extrair_texto_pptx(caminho_pptx):
    print("Lendo slides...")
    prs = Presentation(caminho_pptx)
    texto = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texto += shape.text + "\n"
    return texto

def extrair_texto_site(url):
    print(f"Lendo site: {url}")
    try:
        resp = requests.get(url, timeout=15,
            headers={"User-Agent": "Mozilla/5.0"})
        from html.parser import HTMLParser
        class Parser(HTMLParser):
            def __init__(self):
                super().__init__()
                self.texto = []
                self.ignorar = False
            def handle_starttag(self, tag, attrs):
                if tag in ["script", "style", "nav", "footer"]:
                    self.ignorar = True
            def handle_endtag(self, tag):
                if tag in ["script", "style", "nav", "footer"]:
                    self.ignorar = False
            def handle_data(self, data):
                if not self.ignorar and data.strip():
                    self.texto.append(data.strip())
        p = Parser()
        p.feed(resp.text)
        return " ".join(p.texto)[:5000]
    except:
        return ""

def extrair_texto_foto(caminho_imagem):
    print("Lendo texto da foto com IA...")
    try:
        import base64
        with open(caminho_imagem, "rb") as f:
            img_base64 = base64.b64encode(f.read()).decode()
        resposta = groq_client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            messages=[{
                "role": "user",
                "content": [
                    {"type": "image_url",
                     "image_url": {"url": f"data:image/jpeg;base64,{img_base64}"}},
                    {"type": "text",
                     "text": "Extraia todo o texto visível nesta imagem, mantendo a estrutura original."}
                ]
            }],
            max_tokens=2000
        )
        return resposta.choices[0].message.content
    except Exception as e:
        print(f"Erro ao ler foto: {e}")
        return ""

def coletar_materiais(materiais_apoio):
    if not materiais_apoio:
        return ""
    texto = ""
    for material in materiais_apoio:
        if isinstance(material, str) and material.startswith("INSTRUCAO:"):
            texto += material + "\n\n"
        elif isinstance(material, str) and material.startswith("http"):
            texto += extrair_texto_site(material) + "\n\n"
        elif isinstance(material, str) and material.endswith(".pdf"):
            texto += extrair_texto_pdf(material) + "\n\n"
        elif isinstance(material, str) and material.endswith(".pptx"):
            texto += extrair_texto_pptx(material) + "\n\n"
        elif isinstance(material, str) and material.endswith((".jpg", ".jpeg", ".png")):
            texto += extrair_texto_foto(material) + "\n\n"
    return texto

def buscar_fontes_confiaveis(tema):
    print("Buscando fontes confiáveis em inglês...")
    resultados_todos = []
    try:
        with DDGS() as ddgs:
            for query_template in QUERIES_FONTES[:3]:
                query = query_template.replace("{tema}", tema)
                resultados = list(ddgs.text(query, max_results=2))
                for r in resultados:
                    resultados_todos.append(
                        f"Source: {r.get('href','')}\n{r.get('body','')}"
                    )
                time.sleep(1)
    except Exception as e:
        print(f"Erro na busca: {e}")
    return "\n\n".join(resultados_todos)[:5000]

def verificar_confiabilidade(conceito, fontes):
    resposta = groq_client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{
            "role": "user",
            "content": f"""Analyze these sources about: {conceito}

SOURCES:
{fontes[:3000]}

Answer:
1. Do at least 2 sources agree on the main concept? (yes/no)
2. What do they agree on? (summarize in 2-3 sentences)
3. What is different or contradictory between sources?
4. Confidence level: GREEN (2+ sources agree), YELLOW (only 1 source), RED (sources contradict)

Be concise and factual."""
        }],
        temperature=0.1,
        max_tokens=500
    )
    analise = resposta.choices[0].message.content
    if "GREEN" in analise.upper():
        indicador = "🟢"
    elif "RED" in analise.upper():
        indicador = "🔴"
    else:
        indicador = "🟡"
    return analise, indicador

def analisar_parte(transcricao_parte, numero, total, material_apoio="", fontes=""):
    print(f"Analisando parte {numero}/{total}...")
    resposta = groq_client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[
            {
                "role": "system",
                "content": """You are an expert IT, UX/UI and Computer Science tutor creating study material.
Analyze this class segment thoroughly.
Use support materials and trusted English sources to enrich the analysis.
Always cite sources when using their information.
Correct any errors from the teacher, explaining what is actually correct.
Add confidence indicators: 🟢 confirmed by 2+ sources | 🟡 found in 1 source | 🔴 uncertain
Respond in Brazilian Portuguese."""
            },
            {
                "role": "user",
                "content": f"""CLASS SEGMENT {numero} OF {total}:
{transcricao_parte}

SUPPORT MATERIAL:
{material_apoio[:2000] if material_apoio else 'None provided'}

TRUSTED SOURCES IN ENGLISH:
{fontes[:2000] if fontes else 'None found'}

Complete analysis:
1. Fix transcription errors and explain the correct version
2. Explain ALL mentioned concepts in depth
3. Compare with trusted English sources (cite them)
4. Identify where the teacher simplified or made mistakes
5. Add practical examples and real context
6. Technical glossary (English term → Portuguese translation + pronunciation)
7. Add confidence indicator 🟢🟡🔴 for each major concept"""
            }
        ],
        temperature=0.2,
        max_tokens=4000
    )
    return resposta.choices[0].message.content

def analisar_completo(transcricao_completa, material_apoio="", fontes=""):
    print("Gerando síntese final completa...")
    partes_transcricao = transcricao_completa.split("[Parte ")
    partes_transcricao = [p for p in partes_transcricao if p.strip()]
    analises = []

    for i, parte in enumerate(partes_transcricao):
        analise = analisar_parte(
            parte, i+1, len(partes_transcricao),
            material_apoio, fontes
        )
        analises.append(f"## Analysis of Part {i+1}\n{analise}")
        print(f"Aguardando 20s para não exceder limite do Groq...")
        time.sleep(20)

    analise_unificada = "\n\n".join(analises)
    print("Generating final unified document...")

    sintese_pt = groq_client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[
            {
                "role": "system",
                "content": """Você é um tutor especialista criando material de estudo definitivo.
Crie um documento final unificado, organizado por tópicos (não por partes da aula).
Inclua apenas conteúdo verificado e confiável.
Marque cada conceito com 🟢🟡🔴 de confiabilidade.
Responda em português brasileiro."""
            },
            {
                "role": "user",
                "content": f"""Com base nestas análises:

{analise_unificada[:6000]}

Crie o documento final em PORTUGUÊS com:
# 📚 MATERIAL DE ESTUDO — [título da aula]

## 1. Resumo da Aula
## 2. Conceitos Principais (com 🟢🟡🔴)
## 3. Erros do Professor Corrigidos
## 4. Exemplos Práticos
## 5. Perguntas e Respostas para Prova
## 6. Glossário Técnico (EN → PT + pronúncia)
## 7. Fontes para Aprofundamento"""
            }
        ],
        temperature=0.2,
        max_tokens=4000
    )
    time.sleep(20)

    sintese_en = groq_client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[
            {
                "role": "system",
                "content": """You are an expert tutor creating the definitive study material.
Create a final unified document organized by topics (not by class parts).
Include only verified, reliable content.
Mark each concept with 🟢🟡🔴 confidence indicators.
Respond in clear, professional English."""
            },
            {
                "role": "user",
                "content": f"""Based on these analyses:

{analise_unificada[:6000]}

Create the final document in ENGLISH with:
# 📚 STUDY MATERIAL — [class title]

## 1. Class Summary
## 2. Main Concepts (with 🟢🟡🔴)
## 3. Teacher's Errors Corrected
## 4. Practical Examples
## 5. Exam Q&A
## 6. Technical Glossary (EN term + pronunciation + PT translation)
## 7. Sources for Further Study"""
            }
        ],
        temperature=0.2,
        max_tokens=4000
    )

    doc_pt = sintese_pt.choices[0].message.content
    doc_en = sintese_en.choices[0].message.content
    return doc_pt, doc_en

async def gerar_audio_async(texto, caminho, voz):
    communicate = edge_tts.Communicate(texto, voz)
    await communicate.save(caminho)

def gerar_audio_documento(texto, caminho, voz="en-US-GuyNeural"):
    import re
    from pydub import AudioSegment as AS

    sentencas = re.split(r'(?<=[.!?])\s+', texto)
    sentencas = [s for s in sentencas if not s.startswith('#') and len(s.strip()) > 5]

    partes = []
    parte_atual = ""
    for s in sentencas:
        if len(parte_atual) + len(s) < 500:
            parte_atual += " " + s
        else:
            if parte_atual:
                partes.append(parte_atual.strip())
            parte_atual = s
    if parte_atual:
        partes.append(parte_atual.strip())

    arquivos_temp = []
    for i, parte in enumerate(partes):
        temp = f"/tmp/audio_doc_{i}.mp3"
        asyncio.run(gerar_audio_async(parte, temp, voz))
        arquivos_temp.append(temp)

    combinado = AS.empty()
    for arq in arquivos_temp:
        combinado += AS.from_mp3(arq)
        os.remove(arq)
    combinado.export(caminho, format="mp3")

def enviar_progresso(mensagem):
    requests.post(f"{TELEGRAM_URL}/sendMessage", json={
        "chat_id": CHAT_ID,
        "text": mensagem
    })

def salvar_no_drive(doc_pt, doc_en, titulo="Aula"):
    data = datetime.now().strftime("%Y-%m-%d_%H-%M")
    pasta_base = f"/home/ubuntu/ariana-tutor/materiais/{data}"
    os.makedirs(pasta_base, exist_ok=True)

    # Texto bilíngue PT + EN
    caminho_texto = f"{pasta_base}/{titulo}_bilingual.md"
    with open(caminho_texto, 'w', encoding='utf-8') as f:
        f.write(f"# {titulo} — {data}\n\n")
        f.write("---\n# 🇧🇷 VERSÃO EM PORTUGUÊS\n\n")
        f.write(doc_pt)
        f.write("\n\n---\n# 🇺🇸 ENGLISH VERSION\n\n")
        f.write(doc_en)

    # Áudio EN — voz do americano
    enviar_progresso("🎙️ Generating English audio...")
    caminho_audio_en = f"{pasta_base}/{titulo}_EN.mp3"
    gerar_audio_documento(doc_en, caminho_audio_en, voz="en-US-GuyNeural")

    enviar_progresso("☁️ Uploading to Google Drive...")
    subprocess.run([
        'rclone', 'copy', pasta_base,
        f'gdrive:Ariana Tutor/Videoaulas/{data}'
    ])
   # Gera link direto via rclone
    resultado = subprocess.run([
        'rclone', 'link',
        f'gdrive:Ariana Tutor/Videoaulas/{data}/{titulo}_bilingual.md'
    ], capture_output=True, text=True)
    
    link = resultado.stdout.strip()
    if link:
        enviar_progresso(f"🔗 Acesse seu material aqui:\n{link}")
    
    print(f"Salvo no Google Drive: {data}")
    return caminho_texto

def processar_youtube(url, materiais_apoio=None):
    try:
        enviar_progresso("⬇️ Baixando áudio do YouTube...")
        audio = baixar_audio_youtube(url)
        if not audio:
            enviar_progresso(f"⚠️ Não foi possível baixar: {url}")
            return None, None

        enviar_progresso("🎙️ Transcrevendo com Whisper...")
        transcricao = transcrever_audio(audio)
        if os.path.exists(audio):
            os.remove(audio)

        enviar_progresso("📚 Coletando materiais de apoio...")
        material_texto = coletar_materiais(materiais_apoio)

        enviar_progresso("🔍 Buscando fontes confiáveis...")
        fontes = buscar_fontes_confiaveis(transcricao[:300])

        enviar_progresso("🧠 Analisando conteúdo (pode demorar alguns minutos)...")
        tentativa_analise = 0
        while tentativa_analise < 10:
            try:
                doc_pt, doc_en = analisar_completo(transcricao, material_texto, fontes)
                break
            except Exception as e:
                erro_str = str(e)
                if "rate_limit_exceeded" in erro_str or "429" in erro_str:
                    import re as re2
                    minutos = re2.search(r'(\d+)m', erro_str)
                    segundos = re2.search(r'(\d+\.\d+)s', erro_str)
                    espera = 60
                    if minutos:
                        espera = int(minutos.group(1)) * 60 + 30
                    elif segundos:
                        espera = int(float(segundos.group(1))) + 30
                    enviar_progresso(f"⏳ Limite atingido. Aguardando {espera//60} min...")
                    time.sleep(espera)
                    tentativa_analise += 1
                else:
                    raise

        titulo = f"Aula_{datetime.now().strftime('%Y%m%d_%H%M')}"
        salvar_no_drive(doc_pt, doc_en, titulo)

        enviar_progresso("✅ Pronto! Material salvo no Google Drive 📁\nPasta: Ariana Tutor/Videoaulas")
        return doc_pt, doc_en

    except Exception as e:
        erro_str = str(e)
        if "rate_limit_exceeded" in erro_str or "429" in erro_str:
            import re
            minutos = re.search(r'(\d+)m', erro_str)
            segundos = re.search(r'(\d+\.\d+)s', erro_str)
            espera = 0
            if minutos:
                espera += int(minutos.group(1)) * 60
            if segundos:
                espera += int(float(segundos.group(1)))
            espera += 30
            enviar_progresso(f"⏳ Limite de tokens atingido. Aguardando {espera//60} minutos e tentando de novo...")
            time.sleep(espera)
        enviar_progresso(f"❌ Erro: {e}")
        raise

def processar_arquivo(caminho_video, materiais_apoio=None):
    transcricao = transcrever_audio(caminho_video)
    material_texto = coletar_materiais(materiais_apoio)
    fontes = buscar_fontes_confiaveis(transcricao[:300])
    doc_pt, doc_en = analisar_completo(transcricao, material_texto, fontes)
    salvar_no_drive(doc_pt, doc_en)
    return doc_pt, doc_en

if __name__ == "__main__":
    print("Sistema de transcrição pronto!")